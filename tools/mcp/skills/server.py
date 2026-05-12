"""
MCP Server: Skills - Habilidades avanzadas para el asistente agentic.

Skills disponibles:
- PDF Generator (reportlab)
- Excel Generator (openpyxl)
- PowerPoint Generator (python-pptx)
- Chart Generator (matplotlib)
- Data Analysis (pandas)
- Email Sender (yagmail)
- Web Search (ddgs)
- GitHub Integration (PyGithub)
- GitLab Integration (python-gitlab)
- Slack Messenger (slack-sdk)
- Notion Integration (notion-client)
"""

import os, json, asyncio
from typing import Any
from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import Tool, TextContent

server = Server("cowork-skills")

# ─── Lista de herramientas ────────────────────────────────────

@server.list_tools()
async def list_tools() -> list[Tool]:
    return [
        # Documentos
        Tool(
            name="generate_pdf",
            description="Genera un archivo PDF profesional con título y contenido formateado",
            inputSchema={
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Título del documento PDF"},
                    "content": {"type": "string", "description": "Contenido del PDF (texto o markdown)"},
                    "filename": {"type": "string", "description": "Nombre del archivo sin extensión"},
                    "author": {"type": "string", "description": "Autor del documento (opcional)"}
                },
                "required": ["title", "content", "filename"]
            }
        ),
        Tool(
            name="generate_excel",
            description="Genera un archivo Excel (.xlsx) con datos estructurados y formato profesional",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "Nombre del archivo sin extensión"},
                    "data_json": {"type": "string", "description": "Datos en formato JSON (lista de objetos con claves/valores)"},
                    "sheet_name": {"type": "string", "description": "Nombre de la hoja (opcional)"}
                },
                "required": ["filename", "data_json"]
            }
        ),
        Tool(
            name="generate_pptx",
            description="Genera una presentación PowerPoint (.pptx) con diapositivas",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "Nombre del archivo sin extensión"},
                    "title": {"type": "string", "description": "Título de la presentación"},
                    "slides_json": {"type": "string", "description": "JSON con slides: [{'title':'...', 'content':['bullet1','bullet2']}]"},
                    "theme": {"type": "string", "description": "Tema: default, dark, blue, green (opcional)"}
                },
                "required": ["filename", "title", "slides_json"]
            }
        ),
        # Análisis de datos
        Tool(
            name="generate_chart",
            description="Genera un gráfico (línea, barra, torta) como imagen PNG",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "Nombre del archivo sin extensión"},
                    "chart_type": {"type": "string", "description": "Tipo: line, bar, pie, scatter"},
                    "data_json": {"type": "string", "description": "Datos en JSON: {'x':[valores], 'y':[valores], 'labels':[etiquetas]}"},
                    "title": {"type": "string", "description": "Título del gráfico"}
                },
                "required": ["filename", "chart_type", "data_json", "title"]
            }
        ),
        Tool(
            name="analyze_data",
            description="Analiza datos con pandas y devuelve estadísticas descriptivas",
            inputSchema={
                "type": "object",
                "properties": {
                    "data_json": {"type": "string", "description": "Datos en JSON (lista de objetos)"},
                    "analysis_type": {"type": "string", "description": "Tipo: describe, correlations, groupby"}
                },
                "required": ["data_json"]
            }
        ),
        # Comunicación
        Tool(
            name="send_email",
            description="Envía un email usando Gmail (requiere credenciales configuradas)",
            inputSchema={
                "type": "object",
                "properties": {
                    "to": {"type": "string", "description": "Destinatario(s) separados por coma"},
                    "subject": {"type": "string", "description": "Asunto del email"},
                    "body": {"type": "string", "description": "Cuerpo del mensaje"},
                    "attachment": {"type": "string", "description": "Ruta del archivo adjunto (opcional)"}
                },
                "required": ["to", "subject", "body"]
            }
        ),
        # Búsqueda
        Tool(
            name="web_search",
            description="Busca información actualizada en internet usando DuckDuckGo",
            inputSchema={
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Términos de búsqueda"},
                    "max_results": {"type": "integer", "description": "Máximo de resultados (1-10)", "default": 5}
                },
                "required": ["query"]
            }
        ),
        # Integraciones
        Tool(
            name="github_operation",
            description="Operación con GitHub: list_repos, get_repo, search_code, list_issues",
            inputSchema={
                "type": "object",
                "properties": {
                    "operation": {"type": "string", "description": "list_repos, get_repo, search_code, list_issues"},
                    "query": {"type": "string", "description": "Búsqueda o nombre del repositorio"},
                    "token": {"type": "string", "description": "GitHub token (opcional, usa env GITHUB_TOKEN)"}
                },
                "required": ["operation"]
            }
        ),
        Tool(
            name="slack_message",
            description="Envía un mensaje a un canal de Slack",
            inputSchema={
                "type": "object",
                "properties": {
                    "channel": {"type": "string", "description": "Canal (#general) o usuario (@username)"},
                    "message": {"type": "string", "description": "Texto del mensaje"},
                    "token": {"type": "string", "description": "Slack Bot Token (opcional, usa env SLACK_TOKEN)"}
                },
                "required": ["channel", "message"]
            }
        ),
    ]


# ─── Implementación de cada skill ─────────────────────────────

@server.call_tool()
async def call_tool(name: str, args: dict[str, Any]) -> list[TextContent]:
    
    # ─── PDF ──────────────────────────────────────────────────
    if name == "generate_pdf":
        title = args.get("title", "Documento")
        content = args.get("content", "")
        filename = args.get("filename", "output")
        author = args.get("author", "Cowork-Local AI")
        
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.lib.enums import TA_LEFT, TA_CENTER
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
            
            output_dir = "/media/SSD1T/cowork-local/data/projects"
            os.makedirs(output_dir, exist_ok=True)
            filepath = os.path.join(output_dir, f"{filename}.pdf")
            
            doc = SimpleDocTemplate(filepath, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # Portada
            title_style = ParagraphStyle('CustomTitle', fontSize=24, spaceAfter=30, alignment=TA_CENTER)
            story.append(Paragraph(title, title_style))
            story.append(HRFlowable(width="80%", thickness=2, color="#1a73e8"))
            story.append(Spacer(1, 20))
            
            # Contenido
            body_style = ParagraphStyle('Body', fontSize=11, leading=14, spaceAfter=8)
            for para in content.split('\n\n'):
                if para.strip():
                    clean = para.replace('&','&amp;').replace('<','&lt;').replace('>','&gt;').replace('\n', '<br/>')
                    story.append(Paragraph(clean, body_style))
            
            story.append(Spacer(1, 30))
            story.append(HRFlowable(width="60%", thickness=1))
            story.append(Paragraph(f"Generado por Cowork-Local AI | {author}", ParagraphStyle('Footer', fontSize=8, textColor='#888888')))
            
            doc.build(story)
            return [TextContent(type="text", text=f"✅ PDF generado: {filepath}")]
        except Exception as e:
            return [TextContent(type="text", text=f"❌ Error PDF: {e}")]
    
    # ─── Excel ────────────────────────────────────────────────
    elif name == "generate_excel":
        filename = args.get("filename", "output")
        data_json = args.get("data_json", "[]")
        sheet_name = args.get("sheet_name", "Datos")
        
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            
            data = json.loads(data_json)
            output_dir = "/media/SSD1T/cowork-local/data/projects"
            os.makedirs(output_dir, exist_ok=True)
            filepath = os.path.join(output_dir, f"{filename}.xlsx")
            
            wb = Workbook()
            ws = wb.active
            ws.title = sheet_name
            
            if data:
                headers = list(data[0].keys())
                
                # Estilos
                header_fill = PatternFill(start_color="1a73e8", end_color="1a73e8", fill_type="solid")
                header_font = Font(bold=True, color="FFFFFF", size=11)
                cell_font = Font(size=10)
                thin_border = Border(
                    left=Side(style='thin'), right=Side(style='thin'),
                    top=Side(style='thin'), bottom=Side(style='thin')
                )
                
                # Encabezados
                for col, header in enumerate(headers, 1):
                    cell = ws.cell(row=1, column=col, value=header)
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = Alignment(horizontal='center')
                    cell.border = thin_border
                
                # Datos
                for row_idx, row_data in enumerate(data, 2):
                    for col_idx, header in enumerate(headers, 1):
                        value = row_data.get(header, "")
                        cell = ws.cell(row=row_idx, column=col_idx, value=value)
                        cell.font = cell_font
                        cell.border = thin_border
                
                # Auto-ajustar columnas
                for col in ws.columns:
                    max_length = 0
                    column = col[0].column_letter
                    for cell in col:
                        if cell.value:
                            max_length = max(max_length, len(str(cell.value)))
                    ws.column_dimensions[column].width = min(max_length + 5, 50)
            
            wb.save(filepath)
            return [TextContent(type="text", text=f"✅ Excel generado: {filepath} ({len(data)} filas, {len(headers)} columnas)")]
        except Exception as e:
            return [TextContent(type="text", text=f"❌ Error Excel: {e}")]
    
    # ─── PowerPoint ───────────────────────────────────────────
    elif name == "generate_pptx":
        filename = args.get("filename", "presentacion")
        title = args.get("title", "Presentación")
        slides_json = args.get("slides_json", "[]")
        theme = args.get("theme", "default")
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            
            slides = json.loads(slides_json)
            output_dir = "/media/SSD1T/cowork-local/data/projects"
            os.makedirs(output_dir, exist_ok=True)
            filepath = os.path.join(output_dir, f"{filename}.pptx")
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Colores por tema
            themes = {
                "default": {"bg": RGBColor(255,255,255), "title": RGBColor(0,0,0), "text": RGBColor(51,51,51)},
                "dark": {"bg": RGBColor(30,30,30), "title": RGBColor(255,255,255), "text": RGBColor(200,200,200)},
                "blue": {"bg": RGBColor(240,245,255), "title": RGBColor(26,115,232), "text": RGBColor(51,51,51)},
                "green": {"bg": RGBColor(240,255,245), "title": RGBColor(13,144,79), "text": RGBColor(51,51,51)},
            }
            theme_colors = themes.get(theme, themes["default"])
            
            # Slide de título
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = theme_colors["bg"]
            
            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = theme_colors["title"]
            p.alignment = PP_ALIGN.CENTER
            
            # Slides de contenido
            for slide_data in slides:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                fill.solid()
                fill.fore_color.rgb = theme_colors["bg"]
                
                # Título
                title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = slide_data.get("title", "")
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.color.rgb = theme_colors["title"]
                
                # Bullets
                bullets = slide_data.get("content", [])
                body_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(4.5))
                tf = body_box.text_frame
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {bullet}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = theme_colors["text"]
                    p.space_after = Pt(12)
            
            prs.save(filepath)
            return [TextContent(type="text", text=f"✅ Presentación generada: {filepath} ({len(slides)} diapositivas)")]
        except Exception as e:
            return [TextContent(type="text", text=f"❌ Error PPTX: {e}")]
    
    # ─── Chart ────────────────────────────────────────────────
    elif name == "generate_chart":
        filename = args.get("filename", "grafico")
        chart_type = args.get("chart_type", "bar")
        data_json = args.get("data_json", "{}")
        title = args.get("title", "Gráfico")
        
        try:
            import matplotlib
            matplotlib.use('Agg')
            import matplotlib.pyplot as plt
            
            data = json.loads(data_json)
            output_dir = "/media/SSD1T/cowork-local/data/projects"
            os.makedirs(output_dir, exist_ok=True)
            filepath = os.path.join(output_dir, f"{filename}.png")
            
            fig, ax = plt.subplots(figsize=(10, 6))
            
            if chart_type == "bar":
                ax.bar(data.get("labels", []), data.get("y", []), color='#1a73e8')
            elif chart_type == "line":
                ax.plot(data.get("x", []), data.get("y", []), marker='o', color='#1a73e8')
            elif chart_type == "pie":
                ax.pie(data.get("y", []), labels=data.get("labels", []), autopct='%1.1f%%')
            elif chart_type == "scatter":
                ax.scatter(data.get("x", []), data.get("y", []), color='#1a73e8')
            
            ax.set_title(title, fontsize=16, fontweight='bold')
            ax.grid(True, alpha=0.3)
            plt.tight_layout()
            plt.savefig(filepath, dpi=150)
            plt.close()
            
            return [TextContent(type="text", text=f"✅ Gráfico generado: {filepath}")]
        except Exception as e:
            return [TextContent(type="text", text=f"❌ Error Chart: {e}")]
    
    # ─── Data Analysis ────────────────────────────────────────
    elif name == "analyze_data":
        data_json = args.get("data_json", "[]")
        analysis_type = args.get("analysis_type", "describe")
        
        try:
            import pandas as pd
            import io
            
            data = json.loads(data_json)
            df = pd.DataFrame(data)
            
            if analysis_type == "describe":
                result = df.describe(include='all').to_string()
            elif analysis_type == "correlations":
                numeric_df = df.select_dtypes(include=['number'])
                if not numeric_df.empty:
                    result = numeric_df.corr().to_string()
                else:
                    result = "No hay columnas numéricas para calcular correlaciones"
            elif analysis_type == "groupby":
                result = df.dtypes.to_string()
            else:
                result = df.head(10).to_string()
            
            return [TextContent(type="text", text=f"📊 Análisis ({analysis_type}):\n\n{result}")]
        except Exception as e:
            return [TextContent(type="text", text=f"❌ Error Análisis: {e}")]
    
    # ─── Email ────────────────────────────────────────────────
    elif name == "send_email":
        to = args.get("to", "")
        subject = args.get("subject", "")
        body = args.get("body", "")
        attachment = args.get("attachment", "")
        
        try:
            import yagmail
            
            user = os.getenv("GMAIL_USER", "")
            password = os.getenv("GMAIL_APP_PASSWORD", "")
            
            if not user or not password:
                return [TextContent(type="text", text="⚠️ Configurá GMAIL_USER y GMAIL_APP_PASSWORD en .env")]
            
            yag = yagmail.SMTP(user, password)
            
            if attachment and os.path.exists(attachment):
                yag.send(to=to, subject=subject, contents=body, attachments=attachment)
            else:
                yag.send(to=to, subject=subject, contents=body)
            
            return [TextContent(type="text", text=f"✅ Email enviado a {to}")]
        except Exception as e:
            return [TextContent(type="text", text=f"❌ Error Email: {e}")]
    
    # ─── Web Search ───────────────────────────────────────────
    elif name == "web_search":
        query = args.get("query", "")
        max_results = args.get("max_results", 5)
        
        try:
            from ddgs import DDGS
            
            results = []
            with DDGS() as ddgs:
                for r in ddgs.text(query, max_results=max_results):
                    results.append(f"🔗 **{r['title']}**\n   {r['href']}\n   {r['body'][:250]}")
            
            output = "\n\n".join(results) if results else "Sin resultados para: " + query
            return [TextContent(type="text", text=output)]
        except Exception as e:
            return [TextContent(type="text", text=f"❌ Error Búsqueda: {e}")]
    
    # ─── GitHub ───────────────────────────────────────────────
    elif name == "github_operation":
        operation = args.get("operation", "")
        query = args.get("query", "")
        token = args.get("token", os.getenv("GITHUB_TOKEN", ""))
        
        try:
            from github import Github
            
            g = Github(token) if token else Github()
            
            if operation == "list_repos":
                repos = g.search_repositories(query=query)[:5] if query else g.get_user().get_repos()[:5]
                result = "\n".join([f"📁 {r.full_name} ⭐{r.stargazers_count}" for r in repos])
            elif operation == "search_code":
                repos = g.search_code(query=query)[:5]
                result = "\n".join([f"📄 {r.repository.full_name}: {r.path}" for r in repos])
            elif operation == "list_issues":
                repo = g.get_repo(query)
                issues = repo.get_issues(state='open')[:5]
                result = "\n".join([f"🐛 #{i.number}: {i.title}" for i in issues])
            elif operation == "get_repo":
                repo = g.get_repo(query)
                result = f"📁 {repo.full_name}\n   {repo.description}\n   ⭐{repo.stargazers_count} | Forks: {repo.forks_count}"
            else:
                result = f"Operación: {operation}"
            
            return [TextContent(type="text", text=result)]
        except Exception as e:
            return [TextContent(type="text", text=f"❌ Error GitHub: {e}")]
    
    # ─── Slack ────────────────────────────────────────────────
    elif name == "slack_message":
        channel = args.get("channel", "")
        message = args.get("message", "")
        token = args.get("token", os.getenv("SLACK_TOKEN", ""))
        
        try:
            from slack_sdk import WebClient
            
            if not token:
                return [TextContent(type="text", text="⚠️ Configurá SLACK_TOKEN en .env")]
            
            client = WebClient(token=token)
            response = client.chat_postMessage(channel=channel, text=message)
            
            return [TextContent(type="text", text=f"✅ Mensaje enviado a {channel}")]
        except Exception as e:
            return [TextContent(type="text", text=f"❌ Error Slack: {e}")]
    
    
    if name == "generate_tests":
        code = args.get("code", "")
        import httpx
        import os; key = os.getenv("DEEPSEEK_API_KEY",""); r = httpx.post("https://api.deepseek.com/v1/chat/completions", json={"model":"deepseek-chat","messages":[{"role":"user","content":f"Genera SOLO tests unitarios con pytest:\n\n{code}"}],"max_tokens":2048,"temperature":0.2}, headers={"Authorization":f"Bearer {key}"}, timeout=30)
        return [TextContent(type="text", text=r.json().get("response","") or r.json().get("thinking","")[:3000])]
    
    if name == "review_code":
        code = args.get("code", "")
        import httpx
        import os; key = os.getenv("DEEPSEEK_API_KEY",""); r = httpx.post("https://api.deepseek.com/v1/chat/completions", json={"model":"deepseek-chat","messages":[{"role":"user","content":f"Revisa bugs, mejoras y buenas practicas:\n\n{code}"}],"max_tokens":2048,"temperature":0.2}, headers={"Authorization":f"Bearer {key}"}, timeout=30)
        return [TextContent(type="text", text=r.json().get("response","") or r.json().get("thinking","")[:3000])]
    
    if name == "generate_docs":
        code = args.get("code", "")
        import httpx
        import os; key = os.getenv("DEEPSEEK_API_KEY",""); r = httpx.post("https://api.deepseek.com/v1/chat/completions", json={"model":"deepseek-chat","messages":[{"role":"user","content":f"Genera docstrings y documentacion:\n\n{code}"}],"max_tokens":2048,"temperature":0.2}, headers={"Authorization":f"Bearer {key}"}, timeout=30)
        return [TextContent(type="text", text=r.json().get("response","") or r.json().get("thinking","")[:3000])]
    
    return [TextContent(type="text", text=f"❓ Skill no encontrado: {name}")]


# ─── Punto de entrada ─────────────────────────────────────────
async def main():
    async with stdio_server() as (r, w):
        await server.run(r, w, server.create_initialization_options())

if __name__ == "__main__":
    asyncio.run(main())
