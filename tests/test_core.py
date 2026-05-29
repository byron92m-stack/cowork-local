from __future__ import annotations

import io
from pathlib import Path
from unittest.mock import patch

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Inches

from src.powerpoint_creator.core import (
    SUPPORTED_CHART_TYPES,
    add_bullet_slide,
    add_chart_slide,
    add_content_slide,
    add_image_slide,
    add_image_slide_from_bytes,
    add_table_slide,
    add_title_slide,
    create_presentation,
    save_presentation,
)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def prs():
    return create_presentation()


# ---------------------------------------------------------------------------
# create_presentation
# ---------------------------------------------------------------------------

class TestCreatePresentation:
    def test_default_creation(self):
        p = create_presentation()
        assert len(p.slides) == 0

    def test_with_nonexistent_template_raises(self):
        with pytest.raises(FileNotFoundError, match="Template not found"):
            create_presentation("nonexistent/path.pptx")

    def test_with_template(self, tmp_path):
        template = tmp_path / "template.pptx"
        dummy = Presentation()
        dummy.save(str(template))
        p = create_presentation(str(template))
        assert len(p.slides) == 0


# ---------------------------------------------------------------------------
# add_title_slide
# ---------------------------------------------------------------------------

class TestAddTitleSlide:
    def test_adds_one_slide(self, prs):
        add_title_slide(prs, "Hello")
        assert len(prs.slides) == 1

    def test_title_text(self, prs):
        add_title_slide(prs, "Hello World")
        assert prs.slides[0].shapes.title.text == "Hello World"

    def test_subtitle_set(self, prs):
        add_title_slide(prs, "Title", subtitle="Sub here")
        slide = prs.slides[0]
        assert slide.shapes.title.text == "Title"
        placeholders = [ph for ph in slide.placeholders]
        assert len(placeholders) >= 2

    def test_multiple_title_slides(self, prs):
        add_title_slide(prs, "A")
        add_title_slide(prs, "B")
        assert len(prs.slides) == 2


# ---------------------------------------------------------------------------
# add_content_slide
# ---------------------------------------------------------------------------

class TestAddContentSlide:
    def test_single_paragraph(self, prs):
        add_content_slide(prs, "Slide 1", ["Only paragraph"])
        slide = prs.slides[0]
        assert slide.shapes.title.text == "Slide 1"

    def test_multiple_paragraphs(self, prs):
        add_content_slide(prs, "Multi", ["A", "B", "C"])
        slide = prs.slides[0]
        body = slide.placeholders[1].text_frame
        assert len(body.paragraphs) == 3

    def test_custom_font_size(self, prs):
        add_content_slide(prs, "Size", ["Text"], font_size=24)
        slide = prs.slides[0]
        body = slide.placeholders[1].text_frame
        assert body.paragraphs[0].font.size.pt == 24  # type: ignore[union-attr]

    def test_custom_font_color(self, prs):
        add_content_slide(prs, "Color", ["Text"], font_color=(255, 0, 0))
        slide = prs.slides[0]
        body = slide.placeholders[1].text_frame
        rgb = body.paragraphs[0].font.color.rgb
        assert str(rgb) == "FF0000"


# ---------------------------------------------------------------------------
# add_bullet_slide
# ---------------------------------------------------------------------------

class TestAddBulletSlide:
    def test_bullets_created(self, prs):
        add_bullet_slide(prs, "Bullets", ["One", "Two", "Three"])
        slide = prs.slides[0]
        assert slide.shapes.title.text == "Bullets"

    def test_empty_bullets(self, prs):
        add_bullet_slide(prs, "Empty", [])
        assert len(prs.slides) == 1


# ---------------------------------------------------------------------------
# add_image_slide
# ---------------------------------------------------------------------------

class TestAddImageSlide:
    def test_nonexistent_image_raises(self, prs):
        with pytest.raises(FileNotFoundError, match="Image not found"):
            add_image_slide(prs, "Img", "no-such-file.png")

    def test_valid_image_adds_shape(self, prs, tmp_path):
        from PIL import Image
        img_path = tmp_path / "test.png"
        Image.new("RGB", (10, 10), color="red").save(str(img_path))
        add_image_slide(prs, "Pic", str(img_path))
        slide = prs.slides[0]
        shapes_with_image = [s for s in slide.shapes if hasattr(s, "image")]
        assert len(shapes_with_image) >= 1


# ---------------------------------------------------------------------------
# add_image_slide_from_bytes
# ---------------------------------------------------------------------------

class TestAddImageSlideFromBytes:
    def test_bytes_adds_shape(self, prs):
        from PIL import Image
        buf = io.BytesIO()
        Image.new("RGB", (10, 10), color="blue").save(buf, format="PNG")
        data = buf.getvalue()
        add_image_slide_from_bytes(prs, "From bytes", data)
        slide = prs.slides[0]
        shapes_with_image = [s for s in slide.shapes if hasattr(s, "image")]
        assert len(shapes_with_image) >= 1


# ---------------------------------------------------------------------------
# add_table_slide
# ---------------------------------------------------------------------------

class TestAddTableSlide:
    def test_table_created(self, prs):
        add_table_slide(prs, "Table", ["A", "B"], [["1", "2"], ["3", "4"]])
        slide = prs.slides[0]
        tables = [s for s in slide.shapes if s.has_table]
        assert len(tables) == 1

    def test_table_dimensions(self, prs):
        add_table_slide(prs, "T", ["X", "Y", "Z"], [["a", "b", "c"]])
        slide = prs.slides[0]
        table = next(s for s in slide.shapes if s.has_table).table
        assert table.cell(0, 0).text == "X"
        assert table.cell(1, 2).text == "c"

    def test_empty_headers(self, prs):
        with pytest.raises(ZeroDivisionError):
            add_table_slide(prs, "No data", [], [])


# ---------------------------------------------------------------------------
# add_chart_slide
# ---------------------------------------------------------------------------

class TestAddChartSlide:
    def test_valid_chart(self, prs):
        add_chart_slide(
            prs, "Chart", "bar",
            ["Q1", "Q2"], {"Sales": [100, 200]},
        )
        slide = prs.slides[0]
        chart_shapes = [s for s in slide.shapes if s.has_chart]
        assert len(chart_shapes) == 1

    def test_invalid_chart_type_raises(self, prs):
        with pytest.raises(ValueError, match="Unsupported chart type"):
            add_chart_slide(prs, "Bad", "3d-pie", ["A"], {"S": [1]})

    def test_all_chart_types(self, prs):
        for ct in SUPPORTED_CHART_TYPES:
            p = create_presentation()
            add_chart_slide(p, ct, ct, ["A"], {"S": [1]})
            chart_shapes = [s for s in p.slides[0].shapes if s.has_chart]
            assert len(chart_shapes) == 1

    def test_chart_with_categories_and_series(self, prs):
        add_chart_slide(
            prs, "Data", "line",
            ["Jan", "Feb", "Mar"],
            {"A": [10, 20, 30], "B": [5, 15, 25]},
        )
        slide = prs.slides[0]
        chart = next(s for s in slide.shapes if s.has_chart).chart
        assert len(chart.series) == 2


# ---------------------------------------------------------------------------
# save_presentation
# ---------------------------------------------------------------------------

class TestSavePresentation:
    def test_saves_to_file(self, prs, tmp_path):
        out = tmp_path / "out.pptx"
        result = save_presentation(prs, str(out))
        assert result == out
        assert out.exists()

    def test_invalid_extension_raises(self, prs):
        with pytest.raises(ValueError, match="must have .pptx extension"):
            save_presentation(prs, "out.pdf")

    def test_nested_directory_created(self, prs, tmp_path):
        out = tmp_path / "nested" / "deep" / "out.pptx"
        save_presentation(prs, str(out))
        assert out.exists()

    def test_bytesio_save(self, prs):
        buf = io.BytesIO()
        prs.save(buf)
        assert buf.tell() > 0
